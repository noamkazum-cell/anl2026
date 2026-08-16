# -*- coding: utf-8 -*-
"""Agent360 poster — professor feedback: bold colors, short content, cloud conclusions, large type."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

W, H = 36.0, 48.0

# Bold but not neon — black / deep blue / standout burgundy
NAVY = RGBColor(0x0B, 0x2C, 0x5E)
BLUE = RGBColor(0x1B, 0x4F, 0x8A)
BLACK = RGBColor(0x11, 0x11, 0x11)
BURGUNDY = RGBColor(0x8B, 0x1E, 0x3F)
TEAL = RGBColor(0x0D, 0x5C, 0x63)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLOUD = RGBColor(0xB8, 0xD9, 0xF2)  # stronger light blue — stands out, still readable
CLOUD_EDGE = RGBColor(0x0B, 0x2C, 0x5E)
SOFT_BLUE = RGBColor(0xE8, 0xF1, 0xFA)
SOFT_BURG = RGBColor(0xF8, 0xE8, 0xEE)
SOFT_TEAL = RGBColor(0xE3, 0xF2, 0xF1)

OUT = r"c:\Users\noamk\Documents\BA\Research\anl2026\example posters\Agent360 Poster FINAL.pptx"


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color, width_pt=2.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def no_line(shape):
    shape.line.fill.background()


def add_box(slide, l, t, w, h, fill=None, line=None, line_w=2.0, rounded=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    if rounded:
        try:
            sh.adjustments[0] = 0.06
        except Exception:
            pass
    if fill is None:
        sh.fill.background()
    else:
        set_fill(sh, fill)
    if line is None:
        no_line(sh)
    else:
        set_line(sh, line, line_w)
    return sh


def add_cloud(slide, l, t, w, h, fill=CLOUD, line=CLOUD_EDGE):
    """Cloud-like callout using PowerPoint cloud shape."""
    sh = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(l), Inches(t), Inches(w), Inches(h))
    set_fill(sh, fill)
    set_line(sh, line, 3.0)
    return sh


def set_paragraph(tf, text_parts, font="Calibri", size=32, align=None, space_after=10):
    if isinstance(text_parts, str):
        paragraphs = [[(text_parts, False, BLACK)]]
    elif text_parts and isinstance(text_parts[0], (str, tuple)):
        if all(isinstance(x, str) for x in text_parts):
            paragraphs = [[(x, False, BLACK)] for x in text_parts]
        elif all(isinstance(x, tuple) for x in text_parts):
            paragraphs = [text_parts]
        else:
            paragraphs = [[(x, False, BLACK) if isinstance(x, str) else x for x in text_parts]]
    else:
        paragraphs = []
        for p in text_parts:
            if isinstance(p, str):
                paragraphs.append([(p, False, BLACK)])
            else:
                paragraphs.append([(x, False, BLACK) if isinstance(x, str) else x for x in p])

    for p in list(tf.paragraphs):
        p.clear()
    while len(tf.paragraphs) < len(paragraphs):
        tf.add_paragraph()

    for i, parts in enumerate(paragraphs):
        p = tf.paragraphs[i]
        p.clear()
        if align is not None:
            p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for part in parts:
            text, bold, color = part if not isinstance(part, str) else (part, False, BLACK)
            run = p.add_run()
            run.text = text
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color

    for i in range(len(paragraphs), len(tf.paragraphs)):
        tf.paragraphs[i].clear()


def text_box(slide, l, t, w, h, content, font="Calibri", size=32, align=None, space_after=10, valign=MSO_ANCHOR.TOP):
    sh = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = sh.text_frame
    tf.word_wrap = True
    tf._txBody.bodyPr.set(
        "anchor",
        {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(valign, "t"),
    )
    set_paragraph(tf, content, font=font, size=size, align=align, space_after=space_after)
    return sh


def section_title(slide, l, t, w, title, size=58):
    return text_box(
        slide, l, t, w, 1.15,
        [[(title, True, NAVY)]],
        font="Times New Roman", size=size, align=PP_ALIGN.LEFT, space_after=0,
    )


def pill(slide, l, t, w, h, fill, line, lines, label_color, size=34):
    """lines: str or list of short label lines."""
    add_box(slide, l, t, w, h, fill=fill, line=line, line_w=2.5)
    if isinstance(lines, str):
        lines = lines.split("\n")
    content = [[(ln, True, label_color)] for ln in lines]
    return text_box(
        slide, l, t, w, h, content,
        size=size, align=PP_ALIGN.CENTER, space_after=4, valign=MSO_ANCHOR.MIDDLE,
    )


def arrow_right(slide, l, t, w=0.55, h=0.45):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    set_fill(sh, BLUE)
    set_line(sh, BLUE)
    return sh


def arrow_down(slide, cx, t, h=0.45):
    sh = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(cx - 0.22), Inches(t), Inches(0.44), Inches(h))
    set_fill(sh, BLUE)
    set_line(sh, BLUE)
    return sh


def cloud_takeaway(slide, l, t, w, h, text):
    add_cloud(slide, l, t, w, h)
    return text_box(
        slide, l + 0.55, t + 0.45, w - 1.10, h - 0.90,
        [[(text, True, NAVY)]],
        font="Calibri", size=36, align=PP_ALIGN.CENTER, space_after=0, valign=MSO_ANCHOR.MIDDLE,
    )


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Outer frame — navy
    add_box(slide, 0.30, 0.30, 35.40, 47.40, fill=WHITE, line=NAVY, line_w=3.5, rounded=False)

    # Top accent bar
    add_box(slide, 0.30, 0.30, 35.40, 0.55, fill=NAVY, line=NAVY, line_w=0, rounded=False)

    # ===== HEADER =====
    text_box(
        slide, 1.20, 1.05, 33.60, 2.40,
        [[("Agent360: A Concealment-First Bilateral Negotiator for ANL 2026", True, BURGUNDY)]],
        font="Times New Roman", size=78, align=PP_ALIGN.CENTER, space_after=0,
    )
    text_box(
        slide, 1.20, 3.50, 33.60, 0.70,
        [[("Noam Kazum  ·  Omer Shani Steinmetz  ·  Dr. Galit Haim  ·  Dr. Raz Lin", True, BLACK)]],
        font="Times New Roman", size=40, align=PP_ALIGN.CENTER, space_after=0,
    )
    text_box(
        slide, 1.20, 4.20, 33.60, 0.55,
        [[("The College of Management Academic Studies  ·  ANL / ANAC 2026  ·  June 2026", True, BLUE)]],
        font="Calibri", size=30, align=PP_ALIGN.CENTER, space_after=0,
    )

    LC, RC = 1.15, 18.55
    CW = 16.30
    TX = 0.35

    # =====================================================================
    # LEFT: The Goal
    # =====================================================================
    section_title(slide, LC, 5.10, CW, "The Goal", size=64)
    text_box(
        slide, LC, 6.35, CW, 3.55,
        [
            [
                ("ANL score = ", True, BLACK),
                ("Advantage", True, BURGUNDY),
                (" + ", True, BLACK),
                ("Concealing", True, BLUE),
            ],
            [
                ("Advantage", True, BURGUNDY),
                (" — close good deals above reserved value", True, BLACK),
            ],
            [
                ("Concealing", True, BLUE),
                (" — hide true preferences from the opponent's learner", True, BLACK),
            ],
            [
                ("Design principle: ", True, NAVY),
                ("concealment first, then extract the deal", True, BLACK),
            ],
        ],
        size=38, space_after=16,
    )

    # =====================================================================
    # LEFT: Architecture — labels only, no numbering, no sentences
    # =====================================================================
    section_title(slide, LC, 10.20, CW, "Architecture", size=60)
    text_box(
        slide, LC, 11.40, CW, 0.70,
        [[("Bid stream only — no opponent class, no true utility", True, BLACK)]],
        size=32, space_after=0,
    )

    # Horizontal flow of three layers
    bw = 4.60
    y = 12.30
    # Architecture pills — soft tint OK for short labels on white-ish fills that are clearly colored
    pill(slide, LC, y, bw, 1.50, WHITE, BURGUNDY, "Bidding\nPersona", BURGUNDY, 34)
    arrow_right(slide, LC + bw + 0.15, y + 0.52, 0.60, 0.48)
    pill(slide, LC + bw + 0.90, y, bw, 1.50, WHITE, BLUE, "Opponent\nModel", BLUE, 34)
    arrow_right(slide, LC + 2 * bw + 1.05, y + 0.52, 0.60, 0.48)
    pill(slide, LC + 2 * bw + 1.80, y, bw, 1.50, WHITE, TEAL, "Deal\nExtraction", TEAL, 34)

    # =====================================================================
    # LEFT: Bidding Persona — timeline + short phase labels
    # =====================================================================
    section_title(slide, LC, 14.20, CW, "Bidding Persona", size=60)

    # Gradual preference revelation timeline (relative time t →)
    text_box(
        slide, LC, 15.35, CW, 0.55,
        [[("Relative time  t  →", True, NAVY)]],
        size=28, space_after=0,
    )
    tw = (CW - 1.20) / 3.0  # leave gaps for arrows between segments
    tl_y = 15.95
    tl_h = 1.55
    timeline = [
        (BURGUNDY, "Decoy", "false signal"),
        (RGBColor(0x7A, 0x4A, 0x12), "Transition", "gradual reveal"),
        (BLUE, "Closing", "extract deal"),
    ]
    for i, (edge, title, sub) in enumerate(timeline):
        x = LC + i * (tw + 0.60)
        add_box(slide, x, tl_y, tw, tl_h, fill=WHITE, line=edge, line_w=4.0, rounded=True)
        text_box(
            slide, x, tl_y, tw, tl_h,
            [[(title, True, edge)], [(sub, True, BLACK)]],
            size=30, align=PP_ALIGN.CENTER, space_after=4, valign=MSO_ANCHOR.MIDDLE,
        )
        if i < 2:
            arrow_right(slide, x + tw + 0.08, tl_y + 0.55, 0.44, 0.42)

    # Three phase rows — white fill, bold colored border
    phases = [
        (BURGUNDY, "Decoy", "Wrong issue priorities · maximal mismatch · seat-asymmetric"),
        (RGBColor(0x7A, 0x4A, 0x12), "Transition", "Mix decoy + truth · slower reveal when opening · early exit if conceding"),
        (BLUE, "Closing", "Blend our utility + theirs · mode caps · discount bait"),
    ]
    py = 17.75
    for edge, title, body in phases:
        add_box(slide, LC, py, CW, 2.05, fill=WHITE, line=edge, line_w=4.0)
        text_box(
            slide, LC + 0.40, py + 0.15, CW - 0.80, 1.75,
            [
                [(title, True, edge)],
                [(body, True, BLACK)],
            ],
            size=36, space_after=8, valign=MSO_ANCHOR.MIDDLE,
        )
        py += 2.20

    # Seat note
    text_box(
        slide, LC, 24.45, CW, 1.00,
        [[
            ("Seat asymmetry: ", True, BURGUNDY),
            ("opening seat keeps a longer decoy and stricter gates", True, BLACK),
        ]],
        size=32, space_after=0,
    )

    # =====================================================================
    # LEFT: What we rejected (short)
    # =====================================================================
    section_title(slide, LC, 25.70, CW, "What We Rejected", size=60)
    text_box(
        slide, LC, 26.90, CW, 2.40,
        [
            [
                ("Reverse-psychology", True, BURGUNDY),
                (" (truth early) — hurts Concealing vs learners", True, BLACK),
            ],
            [
                ("Mid-game escape accept", True, BURGUNDY),
                (" — took mediocre / deceptive deals", True, BLACK),
            ],
        ],
        size=34, space_after=14,
    )

    # =====================================================================
    # RIGHT: Opponent Model
    # =====================================================================
    section_title(slide, RC, 5.10, CW, "Opponent Model", size=60)
    text_box(
        slide, RC, 6.30, CW, 4.20,
        [
            [
                ("Smith frequency", True, BLUE),
                (" — base preference estimate from their bids", True, BLACK),
            ],
            [
                ("Recency · time · issue blends", True, BLUE),
                (" — trust late and diverse signals more", True, BLACK),
            ],
            [
                ("Trajectory", True, BLUE),
                (" — concession slope · bait jumps · predicted path", True, BLACK),
            ],
            [
                ("Five modes: ", True, NAVY),
                ("mirror · learner · deceptive · conceding · unknown", True, BLACK),
            ],
        ],
        size=34, space_after=14,
    )

    # =====================================================================
    # RIGHT: Mode response — compact, no gray
    # =====================================================================
    section_title(slide, RC, 10.80, CW, "Mode Response", size=60)

    mode_rows = [
        (TEAL, "Conceding", "Early decoy exit · chase agreement"),
        (BLUE, "Learner", "Full blend · no bait discount"),
        (BURGUNDY, "Deceptive", "Discount Smith · reject bait"),
        (NAVY, "Mirror", "Plain Smith · anti-mirror bids"),
        (TEAL, "Unknown", "Conservative full blend"),
    ]
    my = 12.00
    for edge, name, resp in mode_rows:
        add_box(slide, RC, my, 4.40, 1.15, fill=WHITE, line=edge, line_w=3.5)
        text_box(
            slide, RC, my, 4.40, 1.15,
            [[(name, True, edge)]],
            size=32, align=PP_ALIGN.CENTER, space_after=0, valign=MSO_ANCHOR.MIDDLE,
        )
        text_box(
            slide, RC + 4.60, my, CW - 4.60, 1.15,
            [[(resp, True, BLACK)]],
            size=32, space_after=0, valign=MSO_ANCHOR.MIDDLE,
        )
        my += 1.30

    # =====================================================================
    # RIGHT: Acceptance — short
    # =====================================================================
    section_title(slide, RC, 18.70, CW, "Acceptance", size=60)
    text_box(
        slide, RC, 19.90, CW, 4.00,
        [
            [("Catastrophe floor near deadline", True, BLACK)],
            [("Mode aspiration — slope follows opponent mode", True, BLACK)],
            [("AC-next — as good as our next concealing bid", True, BLACK)],
            [("Bait guard — block deceptive utility spikes", True, BURGUNDY)],
        ],
        size=34, space_after=12,
    )

    # =====================================================================
    # RIGHT: Evaluation — short
    # =====================================================================
    section_title(slide, RC, 24.20, CW, "Evaluation", size=60)
    text_box(
        slide, RC, 25.40, CW, 3.80,
        [
            [
                ("Learners", True, BLUE),
                ("  BOA · MAP · MiCRO", True, BLACK),
            ],
            [
                ("Stress", True, BLUE),
                ("  time-based NegMAS baselines", True, BLACK),
            ],
            [
                ("Sparring", True, BLUE),
                ("  mirror · bait-switch · strong learners", True, BLACK),
            ],
            [
                ("Strong Concealing vs stress · balanced vs learners & deceptive", True, NAVY),
            ],
        ],
        size=34, space_after=14,
    )

    # =====================================================================
    # BOTTOM: Conclusions in CLOUDS — enlarged, very short, stand out
    # =====================================================================
    add_box(slide, 1.00, 29.80, 34.00, 0.22, fill=NAVY, line=NAVY, line_w=0, rounded=False)
    text_box(
        slide, 1.15, 30.20, 33.70, 1.25,
        [[("Conclusions", True, NAVY)]],
        font="Times New Roman", size=72, align=PP_ALIGN.CENTER, space_after=0,
    )

    clouds = [
        "Conceal and extract\nneed separate controls",
        "Test learners, stress\nand deceptive panels",
        "Seat asymmetry\nimproves learner matchups",
        "Decoy early → reveal late\nno privileged info",
    ]
    # 2x2 large clouds — dominant visual
    cw, ch = 15.90, 6.55
    positions = [
        (1.35, 31.60),
        (18.75, 31.60),
        (1.35, 38.55),
        (18.75, 38.55),
    ]
    for (x, y), msg in zip(positions, clouds):
        add_cloud(slide, x, y, cw, ch)
        text_box(
            slide, x + 1.00, y + 1.20, cw - 2.00, ch - 2.40,
            [[(msg, True, NAVY)]],
            font="Calibri", size=48, align=PP_ALIGN.CENTER, space_after=0, valign=MSO_ANCHOR.MIDDLE,
        )

    # Footer
    text_box(
        slide, 1.20, 45.50, 33.60, 1.50,
        [[(
            "ANAC 2026 / ANL  ·  NegMAS SAO  ·  Agent360 final report, College of Management, June 2026",
            True, BLUE
        )]],
        size=26, align=PP_ALIGN.CENTER, space_after=0,
    )

    prs.save(OUT)
    print("Wrote", OUT)


if __name__ == "__main__":
    build()
